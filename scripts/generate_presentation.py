"""
Generates the official REVORA PowerPoint Presentation (.pptx) file.

Styled with modern dark fintech aesthetics (#080c14 base, neon emerald, cyan, violet, rose accents).
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    DARK_BG = RGBColor(8, 12, 20)        # #080c14
    PANEL_BG = RGBColor(15, 23, 42)      # #0f172a
    CARD_BG = RGBColor(30, 41, 59)       # #1e293b
    TEXT_WHITE = RGBColor(248, 250, 252)  # #f8fafc
    TEXT_MUTED = RGBColor(148, 163, 184) # #94a3b8
    EMERALD = RGBColor(16, 185, 129)     # #10b981
    CYAN = RGBColor(6, 182, 212)        # #06b6d4
    VIOLET = RGBColor(139, 92, 246)     # #8b5cf6
    ROSE = RGBColor(244, 63, 94)        # #f43f5e
    GOLD = RGBColor(245, 158, 11)       # #f59e0b

    def add_blank_slide(notes_text=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Background fill
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = DARK_BG
        bg.line.fill.background()
        
        if notes_text:
            slide.notes_slide.notes_text_frame.text = notes_text
        return slide

    def add_header(slide, title, subtitle=None):
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.font.name = "Outfit"
        
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(14)
            p2.font.color.rgb = CYAN
            p2.font.name = "Inter"

    # =========================================================================
    # SLIDE 1: TITLE
    # =========================================================================
    slide1 = add_blank_slide(
        "Speaker Notes:\nGood morning/afternoon judges. I am Shudhanshu Yadav, and today I am excited to present REVORA — a resilient and explainable transaction recovery engine built for the Razorpay RIFT Buildathon 2025. REVORA solves the multi-billion dollar payment failure recovery problem by combining calibrated ML predictions with strict deterministic safety guardrails."
    )
    
    # Title box
    tb = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "REVORA"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = EMERALD
    p.font.name = "Outfit"

    p2 = tf.add_paragraph()
    p2.text = "Resilient & Explainable Transaction Recovery Engine"
    p2.font.size = Pt(26)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_WHITE
    p2.font.name = "Outfit"
    p2.space_before = Pt(10)

    p3 = tf.add_paragraph()
    p3.text = "Razorpay RIFT Buildathon 2025  |  https://razorpay.com/buildathon/"
    p3.font.size = Pt(16)
    p3.font.color.rgb = CYAN
    p3.font.name = "Inter"
    p3.space_before = Pt(15)

    # Presenter Card
    card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.8), Inches(11.333), Inches(1.8))
    card.fill.solid()
    card.fill.fore_color.rgb = PANEL_BG
    card.line.color.rgb = VIOLET

    tf_c = card.text_frame
    tf_c.word_wrap = True
    p_c = tf_c.paragraphs[0]
    p_c.text = "PRESENTER: Shudhanshu Yadav  |  ML/AI Engineer & Data Scientist"
    p_c.font.size = Pt(16)
    p_c.font.bold = True
    p_c.font.color.rgb = TEXT_WHITE
    
    p_c2 = tf_c.add_paragraph()
    p_c2.text = "• MSc Data Science & AI (CUAP)  • BS Data Science (IIT Madras)  • Research Intern (Dhirubhai Ambani Univ)  • AI Intern (Infosys)"
    p_c2.font.size = Pt(13)
    p_c2.font.color.rgb = TEXT_MUTED
    p_c2.space_before = Pt(6)

    p_c3 = tf_c.add_paragraph()
    p_c3.text = "GitHub Repository: https://github.com/skyhitec/REVORA  (58/58 Pytest Passed)"
    p_c3.font.size = Pt(13)
    p_c3.font.color.rgb = EMERALD
    p_c3.space_before = Pt(6)

    # =========================================================================
    # SLIDE 2: THE PROBLEM
    # =========================================================================
    slide2 = add_blank_slide(
        "Speaker Notes:\nPayment failure is not the end of a transaction, but current systems handle it poorly. Blind retries waste fees on dead credentials like expired cards and trigger card network penalties. Meanwhile, retrying flagged transactions risks severe fraud chargebacks. We need a system that answers: Can this transaction be recovered safely without violating security policy?"
    )
    add_header(slide2, "Transaction Failure Is Not the End", "The Problem: Blind Retries, Fraud Risks, & Lack of Auditability")

    # 3 Cards for Problems
    problems = [
        ("1. Blind Retries & Fee Waste", "Indiscriminate retries on dead credentials (EXPIRED_CARD, INVALID_CREDENTIALS) trigger card network fines and waste processing fees.", ROSE),
        ("2. Fraud & Velocity Exposure", "Retrying transactions flagged for fraud or exceeding customer velocity limits leads to high chargeback rates and compliance breaches.", GOLD),
        ("3. Zero Audit Traceability", "Ops and compliance teams cannot verify whether past recovery attempts adhered to risk policies or why retries failed.", CYAN),
    ]
    
    for idx, (title, desc, border_col) in enumerate(problems):
        box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + idx * 3.9), Inches(1.8), Inches(3.7), Inches(3.2))
        box.fill.solid()
        box.fill.fore_color.rgb = PANEL_BG
        box.line.color.rgb = border_col
        
        tf_b = box.text_frame
        tf_b.word_wrap = True
        p_b = tf_b.paragraphs[0]
        p_b.text = title
        p_b.font.size = Pt(16)
        p_b.font.bold = True
        p_b.font.color.rgb = TEXT_WHITE
        
        p_b2 = tf_b.add_paragraph()
        p_b2.text = desc
        p_b2.font.size = Pt(13)
        p_b2.font.color.rgb = TEXT_MUTED
        p_b2.space_before = Pt(12)

    # Core Question Card at bottom
    q_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.3), Inches(11.733), Inches(1.5))
    q_box.fill.solid()
    q_box.fill.fore_color.rgb = CARD_BG
    q_box.line.color.rgb = EMERALD
    
    tf_q = q_box.text_frame
    tf_q.word_wrap = True
    p_q = tf_q.paragraphs[0]
    p_q.text = "THE CORE RECOVERY QUESTION:"
    p_q.font.size = Pt(14)
    p_q.font.bold = True
    p_q.font.color.rgb = EMERALD

    p_q2 = tf_q.add_paragraph()
    p_q2.text = "\"Can the system recover this transaction safely without violating mandatory security & compliance policies?\""
    p_q2.font.size = Pt(18)
    p_q2.font.bold = True
    p_q2.font.color.rgb = TEXT_WHITE
    p_q2.space_before = Pt(6)

    # =========================================================================
    # SLIDE 3: REVORA SOLUTION
    # =========================================================================
    slide3 = add_blank_slide(
        "Speaker Notes:\nREVORA solves this by analyzing transaction context, calculating recovery probability P_rec, evaluating multi-signal risk, running hard safety guardrails, optimizing net Expected Recovery Value, generating human explanations, and recording cryptographically signed audit logs."
    )
    add_header(slide3, "REVORA: Safe Recovery Through Policy-Aware Decisions", "Unified ML Scoring, Deterministic Guardrails, ERV Math, & Cryptographic Audit")

    # Pipeline Diagram
    stages = [
        ("Raw Failure", "Transaction Input", CYAN),
        ("Risk Analysis", "4 Multi-Signal Tiers", VIOLET),
        ("Safety Guardrails", "Hard Policy Override", ROSE),
        ("ERV Calculation", "Net Yield Optimization", EMERALD),
        ("Bounded Action", "7 Decision Actions", GOLD),
        ("SHA-256 Audit", "Append-Only JSONL", CYAN),
    ]

    for idx, (st_name, st_sub, st_color) in enumerate(stages):
        s_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + idx * 1.95), Inches(1.8), Inches(1.8), Inches(1.8))
        s_box.fill.solid()
        s_box.fill.fore_color.rgb = PANEL_BG
        s_box.line.color.rgb = st_color
        
        tf_s = s_box.text_frame
        tf_s.word_wrap = True
        p_s = tf_s.paragraphs[0]
        p_s.text = st_name
        p_s.font.size = Pt(14)
        p_s.font.bold = True
        p_s.font.color.rgb = TEXT_WHITE
        
        p_s2 = tf_s.add_paragraph()
        p_s2.text = st_sub
        p_s2.font.size = Pt(11)
        p_s2.font.color.rgb = TEXT_MUTED
        p_s2.space_before = Pt(8)

    # Feature List Box below
    f_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.9), Inches(11.733), Inches(2.9))
    f_box.fill.solid()
    f_box.fill.fore_color.rgb = PANEL_BG
    f_box.line.color.rgb = EMERALD

    tf_f = f_box.text_frame
    tf_f.word_wrap = True
    
    items = [
        ("Calibrated ML Probability (P_rec)", "XGBoost model predicts exact recovery probability with locked threshold tau* = 0.1600."),
        ("Deterministic Safety Guardrails", "Hard security rules (FRAUD_RISK_BLOCK, non-retryable codes, velocity limits) override ML scores."),
        ("Financial ERV Optimization", "Net ERV = Gross ERV (Amount x P_rec) - Intervention Cost - Friction Penalty."),
        ("Human-Readable Explanations", "Synthesizes clear, plain-language decision rationale per transaction."),
        ("Cryptographic Audit Trail", "Sequential SHA-256 hash chaining (Genesis H_0) provides O(N) linear tamper detection."),
    ]

    for i, (heading, detail) in enumerate(items):
        p_f = tf_f.paragraphs[0] if i == 0 else tf_f.add_paragraph()
        p_f.text = f"• {heading}: {detail}"
        p_f.font.size = Pt(13)
        p_f.font.color.rgb = TEXT_WHITE if i % 2 == 0 else TEXT_MUTED
        p_f.space_before = Pt(8) if i > 0 else Pt(0)

    # =========================================================================
    # SLIDE 4: SYSTEM ARCHITECTURE
    # =========================================================================
    slide4 = add_blank_slide(
        "Speaker Notes:\nHere is REVORA's actual 4-phase micro-architecture. Phase 1 handles clean data generation and leakage checks. Phase 2 trains and serves the calibrated XGBoost model. Phase 3 orchestrates risk scoring, safety guardrails, ERV math, and SHA-256 audit logging. Phase 4 delivers the FastAPI REST backend, transaction simulator, and React Vite web dashboard."
    )
    add_header(slide4, "REVORA Modular System Architecture", "4-Tier Pipeline Built for Reliability, Safety, & Audit Integrity")

    phases = [
        ("Phase 1: Data Engine", "src/data/\n• synthetic_generator.py\n• schema_validator.py\n• leakage_checker.py\n• stratified_splitter.py", CYAN),
        ("Phase 2: ML Engine", "src/ml/\n• feature_pipeline.py\n• model_trainer.py\n• inference_engine.py\n• XGBoost (tau* = 0.1600)", VIOLET),
        ("Phase 3: Policy Engine", "src/decision/ & policy/\n• risk_classifier.py\n• guardrails.py\n• erv_calculator.py\n• engine.py & explainer.py", EMERALD),
        ("Phase 4: Interfaces", "src/api/ & frontend/\n• FastAPI (Port 8000)\n• Stream Simulator\n• React 18 SPA (Port 5173)\n• 58 Pytest Tests Passed", GOLD),
    ]

    for idx, (p_title, p_body, p_col) in enumerate(phases):
        box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + idx * 2.95), Inches(1.8), Inches(2.8), Inches(4.9))
        box.fill.solid()
        box.fill.fore_color.rgb = PANEL_BG
        box.line.color.rgb = p_col
        
        tf_p = box.text_frame
        tf_p.word_wrap = True
        p_t = tf_p.paragraphs[0]
        p_t.text = p_title
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE
        
        p_b = tf_p.add_paragraph()
        p_b.text = p_body
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = TEXT_MUTED
        p_b.space_before = Pt(14)

    # =========================================================================
    # SLIDE 5: SAFETY GUARDRAILS
    # =========================================================================
    slide5 = add_blank_slide(
        "Speaker Notes:\nThis slide demonstrates REVORA's core design philosophy: AI Can Recommend, but Hard Guardrails Decide. In this exact example from our codebase, a transaction has a high ML recovery score of 92%, but because it carries a FRAUD_RISK_BLOCK flag, the mandatory safety guardrail overrides the ML model and forces a BLOCK decision."
    )
    add_header(slide5, "AI Can Recommend. Guardrails Decide.", "Hard Security & Safety Policies Have Higher Authority Than Model Confidence")

    # Guardrail Flow Comparison
    # Left Box: ML Model Recommendation
    left_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = PANEL_BG
    left_box.line.color.rgb = CYAN

    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    p_l = tf_l.paragraphs[0]
    p_l.text = "1. ML Model Scoring (Phase 2)"
    p_l.font.size = Pt(18)
    p_l.font.bold = True
    p_l.font.color.rgb = CYAN

    lines_l = [
        "Transaction ID: demo_fraud_01",
        "Amount: INR 25,000.00",
        "Failure Code: FRAUD_RISK_BLOCK",
        "ML P_rec Probability: 0.9200 (92%)",
        "Threshold Check: 0.92 >= 0.1600 (PASSED)",
        "",
        "ML Recommendation: High Probability Retry",
    ]
    for line in lines_l:
        p_sub = tf_l.add_paragraph()
        p_sub.text = line
        p_sub.font.size = Pt(13)
        p_sub.font.color.rgb = TEXT_WHITE if "High" not in line else EMERALD

    # Right Box: Safety Guardrail Override
    right_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.933), Inches(1.8), Inches(5.6), Inches(4.9))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = PANEL_BG
    right_box.line.color.rgb = ROSE

    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.text = "2. Deterministic Safety Override (Phase 3)"
    p_r.font.size = Pt(18)
    p_r.font.bold = True
    p_r.font.color.rgb = ROSE

    lines_r = [
        "Rule Check: RULE_FRAUD_BLOCK",
        "Condition: failure_code == 'FRAUD_RISK_BLOCK'",
        "Rule Status: FAILED (Mandatory Security Rule)",
        "Forced Decision: BLOCK",
        "Risk Classification: CRITICAL (Score = 100)",
        "",
        "FINAL DECISION: BLOCK (Override Executed)",
    ]
    for line in lines_r:
        p_sub = tf_r.add_paragraph()
        p_sub.text = line
        p_sub.font.size = Pt(13)
        p_sub.font.color.rgb = TEXT_WHITE if "FINAL" not in line else ROSE

    # =========================================================================
    # SLIDE 6: RECOVERY & DEMO SCENARIOS
    # =========================================================================
    slide6 = add_blank_slide(
        "Speaker Notes:\nHere are 4 verified real-world scenarios running live in REVORA's demo studio: 1) Transient Gateway Failure triggers an immediate RETRY; 2) High-Value Auth Failure triggers VIP Escalation; 3) Fraud Risk Flag forces a BLOCK override; and 4) Expired Credit Card triggers a Customer Update Nudge."
    )
    add_header(slide6, "From Failure to Decision: 4 Live Demo Scenarios", "Verified Decision Engine Outputs Across Real-World Failure Scenarios")

    scenarios = [
        ("Transient Gateway Failure", "TEMPORARY_GATEWAY_FAILURE\nAmount: INR 12,500 | P_rec: 85%\nRule Check: PASSED\nAction: RETRY (Net ERV: +10,615)", EMERALD),
        ("VIP High-Value Auth", "AUTHENTICATION_FAILURE\nAmount: INR 45,000 | P_rec: 75%\nHigh-Value Threshold Exceeded\nAction: ESCALATE (Net ERV: +33,725)", VIOLET),
        ("Fraud Risk Override", "FRAUD_RISK_BLOCK\nAmount: INR 25,000 | P_rec: 92%\nRULE_FRAUD_BLOCK Triggered\nAction: BLOCK (Risk: CRITICAL)", ROSE),
        ("Expired Card Credential", "EXPIRED_CARD\nAmount: INR 3,200 | P_rec: 40%\nNon-Retryable Code Rule\nAction: CUSTOMER_ACTION_REQUIRED", GOLD),
    ]

    for idx, (s_title, s_body, s_col) in enumerate(scenarios):
        r_idx = idx // 2
        c_idx = idx % 2
        c_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + c_idx * 5.95), Inches(1.8 + r_idx * 2.5), Inches(5.7), Inches(2.3))
        c_box.fill.solid()
        c_box.fill.fore_color.rgb = PANEL_BG
        c_box.line.color.rgb = s_col
        
        tf_c = c_box.text_frame
        tf_c.word_wrap = True
        p_ct = tf_c.paragraphs[0]
        p_ct.text = s_title
        p_ct.font.size = Pt(16)
        p_ct.font.bold = True
        p_ct.font.color.rgb = s_col
        
        p_cb = tf_c.add_paragraph()
        p_cb.text = s_body
        p_cb.font.size = Pt(12)
        p_cb.font.color.rgb = TEXT_WHITE
        p_cb.space_before = Pt(8)

    # =========================================================================
    # SLIDE 7: EXPLAINABILITY + AUDIT TRAIL
    # =========================================================================
    slide7 = add_blank_slide(
        "Speaker Notes:\nEvery decision made by REVORA leaves cryptographic evidence. Records are formatted into canonical key-sorted JSON and signed using sequential SHA-256 hashing. Starting from our fixed Genesis hash H_0, any record modification or deletion breaks the chain, enabling O(N) linear tamper verification."
    )
    add_header(slide7, "Every Decision Leaves Evidence", "Canonical JSON Serialization & SHA-256 Sequential Hash Chaining")

    # Left: Hash Chain Formulas
    h_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9))
    h_box.fill.solid()
    h_box.fill.fore_color.rgb = PANEL_BG
    h_box.line.color.rgb = CYAN

    tf_h = h_box.text_frame
    tf_h.word_wrap = True
    p_h = tf_h.paragraphs[0]
    p_h.text = "Cryptographic Hash Chain Engine"
    p_h.font.size = Pt(18)
    p_h.font.bold = True
    p_h.font.color.rgb = CYAN

    lines_h = [
        "1. Genesis Hash (H_0):",
        "   H_0 = SHA256('REVORA_PHASE3_GENESIS')",
        "",
        "2. Canonical JSON Serialization:",
        "   C_i = CanonicalJSON(Record_i - {current_hash})",
        "   (Alphabetical keys, compact separators)",
        "",
        "3. Sequential Composite Link (H_i):",
        "   H_i = SHA256(H_{i-1} + C_i)",
        "",
        "4. Tamper Verification:",
        "   AuditVerifier scans JSONL in O(N) linear time.",
    ]
    for line in lines_h:
        p_sub = tf_h.add_paragraph()
        p_sub.text = line
        p_sub.font.size = Pt(12)
        p_sub.font.color.rgb = TEXT_WHITE if ":" in line else TEXT_MUTED

    # Right: Audit Record Schema Fields
    r_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.933), Inches(1.8), Inches(5.6), Inches(4.9))
    r_box.fill.solid()
    r_box.fill.fore_color.rgb = PANEL_BG
    r_box.line.color.rgb = EMERALD

    tf_r = r_box.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.text = "Verified Audit Record Fields"
    p_r.font.size = Pt(18)
    p_r.font.bold = True
    p_r.font.color.rgb = EMERALD

    fields = [
        "• decision_id & transaction_id",
        "• timestamp & policy_version ('1.0.0')",
        "• recovery_probability (P_rec)",
        "• expected_recovery_value & net_expected_recovery_value",
        "• risk_level ('LOW' | 'MEDIUM' | 'HIGH' | 'CRITICAL')",
        "• decision ('RETRY' | 'BLOCK' | 'ESCALATE' | etc.)",
        "• policy_checks (rule_id, passed, forced_decision)",
        "• reason (human-readable explanation string)",
        "• previous_hash & current_hash (SHA-256 hex string)",
    ]
    for f in fields:
        p_f = tf_r.add_paragraph()
        p_f.text = f
        p_f.font.size = Pt(12)
        p_f.font.color.rgb = TEXT_WHITE
        p_f.space_before = Pt(4)

    # =========================================================================
    # SLIDE 8: TESTING & VALIDATION
    # =========================================================================
    slide8 = add_blank_slide(
        "Speaker Notes:\nREVORA is validated by code, not just demonstrated. Our Pytest test suite contains 58 automated tests spanning schema validation, feature engineering, XGBoost calibration, safety guardrails, ERV math, API endpoints, stream processor, and SHA-256 audit verification. All 58 tests pass cleanly."
    )
    add_header(slide8, "Validated, Not Just Demonstrated", "58 Automated Pytest Tests Passed Across All 4 System Phases")

    # Big Metric Stat Box
    m_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(1.5))
    m_box.fill.solid()
    m_box.fill.fore_color.rgb = PANEL_BG
    m_box.line.color.rgb = EMERALD

    tf_m = m_box.text_frame
    tf_m.word_wrap = True
    p_m = tf_m.paragraphs[0]
    p_m.text = "58 PASSED / 0 FAILED / 0 ERRORS"
    p_m.font.size = Pt(32)
    p_m.font.bold = True
    p_m.font.color.rgb = EMERALD

    p_m2 = tf_m.add_paragraph()
    p_m2.text = "Verified Pytest Test Suite Output  |  Execution Time: ~22s  |  (6 non-blocking deprecation warnings)"
    p_m2.font.size = Pt(14)
    p_m2.font.color.rgb = TEXT_WHITE
    p_m2.space_before = Pt(4)

    # 3 Category Cards
    test_cats = [
        ("Phase 1 & 2 Tests (18 Tests)", "• Data Generator & Deterministic Seed\n• Schema Validator & Zero Target Leakage\n• Feature Pipeline & XGBoost Calibration\n• Optimal Threshold (tau* = 0.1600)", CYAN),
        ("Phase 3 Policy Tests (21 Tests)", "• Multi-Signal Risk Scoring Tiers\n• Hard Fraud & Velocity Safety Guardrails\n• Net ERV Financial Math Calculations\n• SHA-256 Audit Chain & Tamper Detection", VIOLET),
        ("Phase 4 System Tests (19 Tests)", "• FastAPI REST API Endpoints (/predict, /decide)\n• Real-Time Stream Processor & Generator\n• Buildathon Presentation Demo Routes\n• Vite Frontend Production Build (npm run build)", GOLD),
    ]

    for idx, (tc_title, tc_body, tc_col) in enumerate(test_cats):
        t_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + idx * 3.9), Inches(3.5), Inches(3.7), Inches(3.2))
        t_box.fill.solid()
        t_box.fill.fore_color.rgb = PANEL_BG
        t_box.line.color.rgb = tc_col
        
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        p_tt = tf_t.paragraphs[0]
        p_tt.text = tc_title
        p_tt.font.size = Pt(15)
        p_tt.font.bold = True
        p_tt.font.color.rgb = TEXT_WHITE
        
        p_tb = tf_t.add_paragraph()
        p_tb.text = tc_body
        p_tb.font.size = Pt(12)
        p_tb.font.color.rgb = TEXT_MUTED
        p_tb.space_before = Pt(10)

    # =========================================================================
    # SLIDE 9: LIVE DEMO / PRODUCT FLOW
    # =========================================================================
    slide9 = add_blank_slide(
        "Speaker Notes:\nDuring live presentation, we can run REVORA in 7 simple steps: Start the server via script, query /health, run a transient retry scenario, run a VIP escalation, demonstrate the fraud block override, inspect the deep-dive modal, and verify the SHA-256 hash chain."
    )
    add_header(slide9, "REVORA in Action: Live Presentation Sequence", "7-Step Product Flow for Buildathon Demonstration")

    demo_steps = [
        ("1. Launch Platform", "Run `scripts\\start_revora.bat` or `python scripts/start_demo.py`"),
        ("2. Verify Service", "Query `GET http://127.0.0.1:8000/health` (Status: OK, Version: 1.0.0)"),
        ("3. Executive Dashboard", "Open `http://localhost:5173` to view Top KPIs & Revenue Recovery Funnel"),
        ("4. Stream Simulator", "Start real-time stream; demonstrate custom event injector with pacing"),
        ("5. Guardrail Override", "Execute `Fraud Risk Flag` scenario -> Observe mandatory BLOCK override"),
        ("6. Transaction Inspector", "Click row to view raw features, ERV math breakdown, & plain-language reason"),
        ("7. Audit Verification", "Click 'Verify Hash Chain' -> Confirm 100% Tamper-Proof SHA-256 chain"),
    ]

    for idx, (step_title, step_cmd) in enumerate(demo_steps):
        s_card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7 + idx * 0.72), Inches(11.733), Inches(0.62))
        s_card.fill.solid()
        s_card.fill.fore_color.rgb = PANEL_BG
        s_card.line.color.rgb = CYAN if idx % 2 == 0 else VIOLET
        
        tf_sc = s_card.text_frame
        tf_sc.word_wrap = True
        p_sc = tf_sc.paragraphs[0]
        p_sc.text = f"{step_title}  —  {step_cmd}"
        p_sc.font.size = Pt(13)
        p_sc.font.bold = True
        p_sc.font.color.rgb = TEXT_WHITE

    # =========================================================================
    # SLIDE 10: IMPACT & FUTURE SCOPE
    # =========================================================================
    slide10 = add_blank_slide(
        "Speaker Notes:\nIn conclusion, REVORA transitions payment recovery from blind retries to intelligent, policy-aware decisioning. It delivers proven ML recovery predictions, strict safety guardrails, net financial optimization, and complete audit transparency. Thank you!"
    )
    add_header(slide10, "From Prototype to Production: Impact & Future Scope", "Transforming Payment Failure Operations with Intelligence & Governance")

    # Left Box: Current Capabilities
    curr_box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.2))
    curr_box.fill.solid()
    curr_box.fill.fore_color.rgb = PANEL_BG
    curr_box.line.color.rgb = EMERALD

    tf_cr = curr_box.text_frame
    tf_cr.word_wrap = True
    p_cr = tf_cr.paragraphs[0]
    p_cr.text = "CURRENT VERIFIED CAPABILITIES"
    p_cr.font.size = Pt(16)
    p_cr.font.bold = True
    p_cr.font.color.rgb = EMERALD

    c_items = [
        "• Calibrated XGBoost recovery predictions (tau* = 0.1600)",
        "• Multi-signal risk engine & hard safety guardrail overrides",
        "• Expected Recovery Value (ERV) net yield optimization",
        "• Append-only SHA-256 tamper-proof audit log & linear verifier",
        "• Production FastAPI REST service & 8 endpoints",
        "• Real-time payment stream simulator with seed replay",
        "• Single-Page React 18 + Vite web dashboard (5 live views)",
    ]
    for ci in c_items:
        p_sub = tf_cr.add_paragraph()
        p_sub.text = ci
        p_sub.font.size = Pt(12)
        p_sub.font.color.rgb = TEXT_WHITE
        p_sub.space_before = Pt(4)

    # Right Box: Future Scope
    fut_box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.933), Inches(1.8), Inches(5.6), Inches(4.2))
    fut_box.fill.solid()
    fut_box.fill.fore_color.rgb = PANEL_BG
    fut_box.line.color.rgb = VIOLET

    tf_fu = fut_box.text_frame
    tf_fu.word_wrap = True
    p_fu = tf_fu.paragraphs[0]
    p_fu.text = "FUTURE EXPANSION ROADMAP"
    p_fu.font.size = Pt(16)
    p_fu.font.bold = True
    p_fu.font.color.rgb = VIOLET

    f_items = [
        "• Multi-Gateway Adaptive Routing across provider stacks",
        "• Contextual Bandit Reinforcement Learning for dynamic fees",
        "• Outbound WebHook Event Subscriptions for merchant ERPs",
        "• Distributed Kafka/EventHub stream ingestion pipeline",
        "• Automated regulatory compliance report exporter (PDF/CSV)",
    ]
    for fi in f_items:
        p_sub = tf_fu.add_paragraph()
        p_sub.text = fi
        p_sub.font.size = Pt(12)
        p_sub.font.color.rgb = TEXT_MUTED
        p_sub.space_before = Pt(6)

    # Final Closing Banner
    banner = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.15), Inches(11.733), Inches(0.85))
    banner.fill.solid()
    banner.fill.fore_color.rgb = CARD_BG
    banner.line.color.rgb = CYAN

    tf_ba = banner.text_frame
    tf_ba.word_wrap = True
    p_ba = tf_ba.paragraphs[0]
    p_ba.text = "REVORA combines recovery yield, safety guardrails, explainability, and auditability into one seamless decision pipeline."
    p_ba.font.size = Pt(14)
    p_ba.font.bold = True
    p_ba.font.color.rgb = CYAN
    p_ba.alignment = PP_ALIGN.CENTER

    # Save output file
    out_path = Path(__file__).resolve().parent.parent / "docs" / "REVORA_Presentation.pptx"
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"[SUCCESS] Generated presentation PowerPoint deck: {out_path}")


if __name__ == "__main__":
    create_deck()
